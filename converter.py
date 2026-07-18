r"""
Docling Project Converter  (MCP-oriented, library module)
-----------------------------------------------------------
Converts every supported file inside a project's `source/` folder to
Markdown (and, where possible, JSON + table CSVs) inside that same
project's `processed/` folder. Designed to be called by an MCP server
(see server.py) with a fixed, predictable folder layout -- no
interactive prompts, no dated output folders.

Expected project layout (folder_path is the project root):

    <folder_path>/
        source/                  <- you put the raw files here
            Annual Report.pdf
            Investor Presentation.pptx
            Financial Statements.xlsx
        processed/                <- created automatically
            markdown/
                Annual Report.md
                Investor Presentation.md
            json/
                Financial Statements.json
            tables/
                Financial Statements__table_01.csv
        metadata.json             <- written after every run; read this first

`source/` is never modified -- files are read in place, not moved or
archived, so it stays a stable location to point Claude (or yourself)
back at the original document. Re-running convert_folder() on a
project that's already been converted only (re)converts files that are
new or whose modified-time has changed since the last run (or every
file, if force=True) -- see metadata.json's "source_mtime" per file.

Legacy Excel formats (.xls, .xlsb, .xlsm) are automatically converted
to .xlsx via Excel (requires Excel installed on Windows) before Docling
processes them.

LARGE-DOCUMENT CHUNKING:
Docling can silently truncate very large documents partway through
(e.g. stopping around page 100-150 of a big PDF) instead of raising an
error. To work around this, files above a page/slide/section threshold
are split into smaller chunks BEFORE being sent to Docling, converted
chunk-by-chunk, and the resulting Markdown is stitched back together
into a single output file. All temporary chunk files (and temp .xlsx
files from legacy Excel conversion) are deleted afterward.

    PDF   -> split by page count (via pypdf)
    PPTX  -> split by slide count (via python-pptx)
    DOCX  -> split by top-level block count: paragraphs + tables
             (via python-docx) -- DOCX has no fixed "page" concept,
             so this is a proxy unit, not a literal page count.

NOTE: a chunked file has no single combined DoclingDocument object
(only a stitched Markdown string), so JSON and table-CSV export only
ever happen for files that convert WHOLE, i.e. at or under
chunk_threshold. A file over the threshold (most 10-Ks, most long
decks) will only produce Markdown. This is a real limitation, not a
bug to chase -- see server.py's module docstring for why that's an
acceptable tradeoff for this project's actual downstream use.

Supported input types:
    PDF, DOCX, PPTX, XLSX, HTML, AsciiDoc,
    PNG, JPG, JPEG, TIFF, BMP, WEBP,
    XLS, XLSB, XLSM (auto-converted to XLSX first)

Entry points:
    convert_folder(folder_path, chunk_threshold=25, force=False) -> dict
        Runs conversion and returns/writes the metadata.json dict.
    read_metadata(folder_path) -> dict | None
        Reads an existing metadata.json without converting anything.

Requirements:
    pip install docling pywin32 pypdf python-docx python-pptx pandas
"""

import copy
import json
import logging
import shutil
import time
from datetime import datetime
from pathlib import Path

from docling.document_converter import DocumentConverter
from docling.datamodel.base_models import ConversionStatus

# All file types Docling supports natively
SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".html",
    ".htm",
    ".adoc",
    ".png",
    ".jpg",
    ".jpeg",
    ".tiff",
    ".tif",
    ".bmp",
    ".webp",
}

# Legacy Excel formats that need pre-conversion to .xlsx
LEGACY_EXCEL_EXTENSIONS = {".xls", ".xlsb", ".xlsm"}

# File types eligible for chunking, and the unit size per chunk.
# (Same chunk size used across types for consistency, per user preference.)
CHUNK_SIZE = 25
CHUNKABLE_EXTENSIONS = {".pdf", ".docx", ".pptx"}

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  %(levelname)s  %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger(__name__)


def convert_legacy_excel_to_xlsx(file: Path, dest_dir: Path) -> Path | None:
    """
    Opens a legacy Excel file via the Excel COM interface and saves it
    as .xlsx inside dest_dir (a temp working folder, NOT the source folder --
    saving next to the source risks silently overwriting and later deleting
    a real .xlsx that happens to share the same stem, e.g. "budget.xls" and
    "budget.xlsx" both present). Returns the new .xlsx Path, or None on failure.
    Requires Excel to be installed and pywin32.
    """
    try:
        import win32com.client  # type: ignore
    except ImportError:
        log.error("  pywin32 not installed. Run: pip install pywin32")
        return None

    xlsx_path = dest_dir / f"{file.stem}.xlsx"
    excel = None
    wb = None
    try:
        log.info(f"  → Opening Excel to convert {file.name} → {xlsx_path.name}")
        excel = win32com.client.Dispatch("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False

        wb = excel.Workbooks.Open(str(file.resolve()))
        # 51 = xlOpenXMLWorkbook (.xlsx)
        wb.SaveAs(str(xlsx_path.resolve()), FileFormat=51)
        wb.Close(False)
        log.info(f"  ✓ Excel conversion done: {xlsx_path.name}")
        return xlsx_path

    except Exception as e:
        log.error(f"  ✗ Excel conversion failed for {file.name}: {e}")
        if wb:
            try:
                wb.Close(False)
            except Exception:
                pass
        return None

    finally:
        if excel:
            try:
                excel.Quit()
            except Exception:
                pass


# ---------------------------------------------------------------------------
# Chunking helpers
# ---------------------------------------------------------------------------

def get_chunk_count(file: Path) -> int:
    """
    Return how many chunk-units (pages/slides/blocks) a file has, used to
    decide whether chunking is needed. Returns 1 if it can't be determined
    (in which case the file is processed whole, as before).
    """
    suffix = file.suffix.lower()
    try:
        if suffix == ".pdf":
            from pypdf import PdfReader
            return len(PdfReader(str(file)).pages)

        if suffix == ".pptx":
            from pptx import Presentation
            return len(Presentation(str(file)).slides)

        if suffix == ".docx":
            from docx import Document
            doc = Document(str(file))
            return len(doc.element.body)  # top-level XML blocks (paras+tables+etc.)

    except Exception as e:
        log.warning(f"  Could not inspect {file.name} to check size: {e}")

    return 1


def split_pdf(file: Path, chunk_dir: Path, chunk_size: int) -> list[Path]:
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(str(file))
    total = len(reader.pages)
    chunks = []
    for start in range(0, total, chunk_size):
        end = min(start + chunk_size, total)
        writer = PdfWriter()
        for p in range(start, end):
            writer.add_page(reader.pages[p])
        chunk_path = chunk_dir / f"{file.stem}__chunk_{start + 1:04d}-{end:04d}.pdf"
        with open(chunk_path, "wb") as f:
            writer.write(f)
        chunks.append(chunk_path)
    return chunks


def split_pptx(file: Path, chunk_dir: Path, chunk_size: int) -> list[Path]:
    """
    Splits a PPTX into chunks of `chunk_size` slides by copying the whole
    deck N times and deleting out-of-range slides from each copy.
    (python-pptx has no native slide-removal API, so we drop to the
    underlying XML to remove slides cleanly.)
    """
    from pptx import Presentation

    prs = Presentation(str(file))
    total = len(prs.slides)
    chunks = []

    xml_slides = prs.slides._sldIdLst  # the <p:sldIdLst> element
    all_slide_ids = list(xml_slides)

    for start in range(0, total, chunk_size):
        end = min(start + chunk_size, total)
        chunk_path = chunk_dir / f"{file.stem}__chunk_{start + 1:04d}-{end:04d}.pptx"
        shutil.copy(file, chunk_path)

        chunk_prs = Presentation(str(chunk_path))
        sld_id_lst = chunk_prs.slides._sldIdLst
        ids_to_remove = list(sld_id_lst)[:start] + list(sld_id_lst)[end:]
        for sld_id in ids_to_remove:
            rId = sld_id.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            chunk_prs.part.drop_rel(rId)
            sld_id_lst.remove(sld_id)

        chunk_prs.save(str(chunk_path))
        chunks.append(chunk_path)

    return chunks


def split_docx(file: Path, chunk_dir: Path, chunk_size: int) -> list[Path]:
    """
    Splits a DOCX into chunks of `chunk_size` top-level body blocks
    (paragraphs/tables/etc). This is a structural proxy for "pages" since
    DOCX has no fixed page boundaries until rendered.
    """
    from docx import Document

    doc = Document(str(file))
    body = doc.element.body
    all_blocks = list(body)
    # last element is usually sectPr (section properties) -- keep it out of
    # the slicing, we'll re-append it to each chunk for valid section formatting
    sect_pr = None
    if len(all_blocks) and all_blocks[-1].tag.endswith("}sectPr"):
        sect_pr = all_blocks[-1]
        all_blocks = all_blocks[:-1]

    total = len(all_blocks)
    chunks = []

    for start in range(0, total, chunk_size):
        end = min(start + chunk_size, total)
        chunk_doc = Document()
        chunk_body = chunk_doc.element.body
        # clear default empty paragraph docx ships with
        for child in list(chunk_body):
            chunk_body.remove(child)

        for block in all_blocks[start:end]:
            chunk_body.append(copy.deepcopy(block))

        if sect_pr is not None:
            chunk_body.append(copy.deepcopy(sect_pr))

        chunk_path = chunk_dir / f"{file.stem}__chunk_{start + 1:04d}-{end:04d}.docx"
        chunk_doc.save(str(chunk_path))
        chunks.append(chunk_path)

    return chunks


def split_file(file: Path, chunk_dir: Path, chunk_size: int) -> list[Path]:
    suffix = file.suffix.lower()
    if suffix == ".pdf":
        return split_pdf(file, chunk_dir, chunk_size)
    if suffix == ".pptx":
        return split_pptx(file, chunk_dir, chunk_size)
    if suffix == ".docx":
        return split_docx(file, chunk_dir, chunk_size)
    raise ValueError(f"Unsupported chunk type: {suffix}")


# ---------------------------------------------------------------------------
# Conversion
# ---------------------------------------------------------------------------

def write_tables(document: object, dest_dir: Path, stem: str) -> list[dict]:
    """
    Exports every table found in a whole (unchunked) DoclingDocument to its
    own CSV file in dest_dir, named "{stem}__table_{NN}.csv".

    This is deliberately generic -- it does NOT try to guess which table is
    "the income statement" vs "the balance sheet" (Docling has no concept of
    that; it only knows "there is a table on page N"). Each entry records the
    source page number, if Docling reported one, so a reader can jump
    straight to the right table without a semantic-labeling pass.

    Returns a list of {"file": "<stem>__table_01.csv", "page": int|None}
    dicts, one per table written. Returns [] if the document has no tables,
    or if a table fails to export (logged as a warning, not fatal).
    """
    tables = getattr(document, "tables", None) or []
    written = []
    for i, table in enumerate(tables, 1):
        csv_name = f"{stem}__table_{i:02d}.csv"
        try:
            df = table.export_to_dataframe(document)
            csv_path = dest_dir / csv_name
            df.to_csv(csv_path, index=False)
            page = None
            if getattr(table, "prov", None):
                page = table.prov[0].page_no
            written.append({"file": csv_name, "page": page})
        except Exception as e:
            log.warning(f"    Could not export table {i} of {stem} to CSV: {e}")
    return written


def convert_one(converter: DocumentConverter, file: Path) -> tuple[object | None, str | None, str, list[str]]:
    """
    Converts a single file with Docling.
    Returns (document_or_None, markdown_text_or_None, status_str, error_messages).
    status_str is one of: "success", "partial", "failed", "exception".
    `document` is the Docling DoclingDocument object (needed for non-Markdown
    exports, e.g. JSON for Excel files) -- None whenever conversion produced
    no usable document.
    """
    try:
        result = converter.convert(file)
        errors = [e.error_message for e in getattr(result, "errors", [])]

        if result.status == ConversionStatus.SUCCESS:
            return result.document, result.document.export_to_markdown(), "success", errors

        elif result.status == ConversionStatus.PARTIAL_SUCCESS:
            return result.document, result.document.export_to_markdown(), "partial", errors

        else:
            return None, None, "failed", errors

    except Exception as e:
        return None, None, "exception", [str(e)]


def convert_with_chunking(
    converter: DocumentConverter,
    file: Path,
    chunk_threshold: int,
    chunk_size: int,
    work_dir: Path,
) -> tuple[object | None, str | None, str, list[str], bool]:
    """
    Decides, per file, where Docling actually reads from:

      NOT SPLIT  -> file is under/at the threshold (or not a chunkable
                    type, or splitting failed). Docling converts the
                    ORIGINAL file, in place, from its source folder.
                    No temp folder is created at all.

      SPLIT      -> file exceeds the threshold. The file is split into
                    chunk files inside a freshly-created temp folder
                    under `work_dir`. Docling converts EACH CHUNK from
                    that temp folder (never the original). The temp
                    folder and its chunk files are deleted once all
                    chunks for this file have been converted.

    Returns (document, markdown, status, errors, was_chunked) -- the extra
    `was_chunked` bool (vs. convert_one's 4-tuple) is what convert_folder()
    records in metadata.json, so it reflects what actually happened even in
    the edge case where splitting itself fails and this falls back to
    NOT SPLIT despite the file being over threshold.
    """
    suffix = file.suffix.lower()

    # --- Decide split vs. no-split -----------------------------------
    if suffix not in CHUNKABLE_EXTENSIONS:
        log.info(f"  → NOT SPLIT ({suffix} is not a chunkable type) — converting original file in place")
        return (*convert_one(converter, file), False)

    unit_count = get_chunk_count(file)
    if unit_count <= chunk_threshold:
        log.info(
            f"  → NOT SPLIT ({unit_count} units <= threshold {chunk_threshold}) — "
            f"converting original file in place: {file}"
        )
        return (*convert_one(converter, file), False)

    # --- Split path: create temp folder, convert chunks from there --
    log.info(
        f"  → SPLIT ({unit_count} units > threshold {chunk_threshold}) — "
        f"creating temp chunk folder and converting from there"
    )

    chunk_dir = work_dir / f"_chunks_{file.stem}"
    chunk_dir.mkdir(parents=True, exist_ok=True)
    log.info(f"    Temp chunk folder: {chunk_dir}")

    try:
        chunk_files = split_file(file, chunk_dir, chunk_size)
    except Exception as e:
        log.warning(
            f"  Could not split {file.name} ({e}); falling back to NOT SPLIT — "
            f"converting original file in place instead"
        )
        shutil.rmtree(chunk_dir, ignore_errors=True)
        return (*convert_one(converter, file), False)

    md_parts = []
    all_errors = []
    overall_status = "success"

    for i, chunk in enumerate(chunk_files, 1):
        log.info(f"    [chunk {i}/{len(chunk_files)}] Converting from temp folder: {chunk}")
        _doc, md, status, errors = convert_one(converter, chunk)
        # Note: each chunk's `document` object is intentionally discarded here.
        # There's no single combined DoclingDocument for a split file -- only
        # a stitched Markdown string -- so JSON export (see main()) is only
        # ever attempted for whole, unchunked files. This is a non-issue for
        # Excel specifically, since .xlsx is never in CHUNKABLE_EXTENSIONS.

        if md:
            md_parts.append(md)
        if status in ("partial", "failed", "exception"):
            overall_status = "partial" if overall_status == "success" else overall_status
        if status in ("failed", "exception"):
            log.warning(f"      ⚠ Chunk {chunk.name} returned status={status}: {errors}")
        all_errors.extend(errors)

    # cleanup temp chunk files -- the temp folder this file's chunks lived in
    # is removed now that every chunk has been converted, win or lose.
    log.info(f"    Cleaning up temp chunk folder: {chunk_dir}")
    shutil.rmtree(chunk_dir, ignore_errors=True)

    if not md_parts:
        return None, None, "failed", all_errors, True

    combined_md = "\n\n<!-- ===== chunk boundary ===== -->\n\n".join(md_parts)
    return None, combined_md, overall_status, all_errors, True


# ---------------------------------------------------------------------------
# MCP-facing entry points
# ---------------------------------------------------------------------------

def read_metadata(folder_path) -> dict | None:
    """
    Reads an existing metadata.json from a project folder WITHOUT converting
    anything. Returns None if the project has no metadata.json yet (i.e. it
    has never been converted).
    """
    metadata_path = Path(folder_path) / "metadata.json"
    if not metadata_path.exists():
        return None
    with metadata_path.open("r", encoding="utf-8") as fp:
        return json.load(fp)


def _is_up_to_date(origin: Path, previous_files: dict, folder_path: Path, force: bool) -> bool:
    """Whether `origin` can be skipped this run because it already converted
    successfully last time and hasn't changed since."""
    if force:
        return False
    prev = previous_files.get(origin.name)
    if prev is None or prev.get("status") not in ("success", "partial"):
        return False
    if prev.get("source_mtime") != origin.stat().st_mtime:
        return False
    md = prev.get("markdown")
    if md and not (folder_path / md).exists():
        return False
    return True


def _rel(path: Path, folder_path: Path) -> str:
    """Path relative to the project root, with forward slashes (for
    metadata.json portability across OSes)."""
    return str(path.relative_to(folder_path)).replace("\\", "/")


def convert_folder(folder_path, chunk_threshold: int = CHUNK_SIZE, force: bool = False) -> dict:
    """
    Converts every supported file in <folder_path>/source/ into
    <folder_path>/processed/{markdown,json,tables}/, and writes
    <folder_path>/metadata.json summarizing the result. See the module
    docstring for the full folder layout and the metadata.json schema.

    folder_path : project root. Must already contain a source/ subfolder.
    chunk_threshold : pages/slides/blocks above which a file is split into
        chunks before conversion (protects against Docling silently
        truncating large documents). Files at or under this size convert
        whole, which is required for JSON/table export.
    force : if True, reconverts every file regardless of prior runs. If
        False (default), a file already recorded as success/partial in
        the previous metadata.json, with an unchanged modified-time and
        an on-disk markdown file, is left alone and carried forward as-is
        -- only new or changed files are (re)converted.

    Returns the metadata dict that was written to metadata.json.
    Raises FileNotFoundError if <folder_path>/source/ doesn't exist.

    source/ itself is never modified -- files are read in place, never
    moved or deleted, so it stays a stable place to point back at the
    original document.
    """
    folder_path = Path(folder_path)
    source_dir = folder_path / "source"
    processed_dir = folder_path / "processed"
    md_dir = processed_dir / "markdown"
    json_dir = processed_dir / "json"
    tables_dir = processed_dir / "tables"
    work_dir = folder_path / ".tmp_conversion"

    if not source_dir.is_dir():
        raise FileNotFoundError(
            f"No source/ folder found at {source_dir}. Create it and add "
            f"the files you want converted, then try again."
        )

    for d in (md_dir, json_dir, tables_dir, work_dir):
        d.mkdir(parents=True, exist_ok=True)

    previous = read_metadata(folder_path) or {}
    previous_files = {f["name"]: f for f in previous.get("files", [])}

    all_files = sorted(source_dir.iterdir())
    native_files = [f for f in all_files if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS]
    legacy_excel = [f for f in all_files if f.is_file() and f.suffix.lower() in LEGACY_EXCEL_EXTENSIONS]
    skipped_unsupported = sorted(
        f.name for f in all_files
        if f.is_file()
        and f.suffix.lower() not in SUPPORTED_EXTENSIONS
        and f.suffix.lower() not in LEGACY_EXCEL_EXTENSIONS
    )
    if skipped_unsupported:
        log.info(f"Skipping {len(skipped_unsupported)} unsupported file(s): {skipped_unsupported}")

    # Split into "already up to date, carry the old entry forward" vs.
    # "needs (re)conversion" -- BEFORE doing any legacy-Excel pre-conversion,
    # so files that haven't changed never trigger an Excel COM launch.
    files_out = []
    native_to_convert = []
    for f in native_files:
        if _is_up_to_date(f, previous_files, folder_path, force):
            files_out.append(previous_files[f.name])
        else:
            native_to_convert.append(f)

    legacy_to_convert = []
    for f in legacy_excel:
        if _is_up_to_date(f, previous_files, folder_path, force):
            files_out.append(previous_files[f.name])
        else:
            legacy_to_convert.append(f)

    # Pre-convert only the legacy Excel files that actually need it.
    origin_of = {f: f for f in native_to_convert}
    converted_xlsx = []
    if legacy_to_convert:
        log.info(f"Found {len(legacy_to_convert)} legacy Excel file(s) to pre-convert: {[f.name for f in legacy_to_convert]}")
        for f in legacy_to_convert:
            xlsx = convert_legacy_excel_to_xlsx(f, work_dir)
            if xlsx:
                converted_xlsx.append(xlsx)
                origin_of[xlsx] = f
            else:
                # Excel conversion itself failed -- record it as a failure
                # rather than silently dropping the file.
                files_out.append({
                    "name": f.name,
                    "type": f.suffix.lower().lstrip("."),
                    "status": "failed",
                    "markdown": None, "json": None, "tables": [],
                    "chunked": False,
                    "source_mtime": f.stat().st_mtime,
                    "converted_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "errors": ["Legacy Excel -> .xlsx pre-conversion failed (see server log)."],
                })

    seen = set()
    to_convert = []
    for f in native_to_convert + converted_xlsx:
        if f.name not in seen:
            seen.add(f.name)
            to_convert.append(f)
    to_convert.sort(key=lambda f: f.name)

    if not to_convert:
        log.info("Nothing new to convert -- every file is already up to date.")
    else:
        log.info(f"Converting {len(to_convert)} file(s) with Docling ({len(previous_files)} previously converted, reused where unchanged)")
        log.info(f"Chunk threshold: {chunk_threshold} pages/slides/blocks")
        converter = DocumentConverter()
        start_time = time.time()

        for i, file in enumerate(to_convert, 1):
            origin = origin_of[file]
            log.info(f"[{i}/{len(to_convert)}] Converting: {origin.name}  ({origin.suffix.upper()})")

            document, md_text, status, errors, chunked = convert_with_chunking(
                converter, file, chunk_threshold, chunk_threshold, work_dir
            )

            entry = {
                "name": origin.name,
                "type": origin.suffix.lower().lstrip("."),
                "status": status,
                "markdown": None,
                "json": None,
                "tables": [],
                "chunked": chunked,
                "source_mtime": origin.stat().st_mtime,
                "converted_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "errors": errors or [],
            }

            if md_text is not None:
                md_path = md_dir / f"{file.stem}.md"
                md_path.write_text(md_text, encoding="utf-8")
                entry["markdown"] = _rel(md_path, folder_path)

            # JSON + table export only happen for a WHOLE document -- a
            # chunked file has no single combined DoclingDocument (see
            # module docstring). document is None whenever conversion
            # failed outright, or the file was split into chunks.
            if document is not None:
                json_path = json_dir / f"{file.stem}.json"
                try:
                    with json_path.open("w", encoding="utf-8") as fp:
                        json.dump(document.export_to_dict(), fp, ensure_ascii=False, indent=2)
                    entry["json"] = _rel(json_path, folder_path)
                except Exception as e:
                    log.warning(f"  Could not write JSON export for {origin.name}: {e}")

                table_entries = write_tables(document, tables_dir, file.stem)
                entry["tables"] = [
                    {"file": _rel(tables_dir / t["file"], folder_path), "page": t["page"]}
                    for t in table_entries
                ]

            if status == "success":
                log.info(f"  \u2713 {origin.name}: markdown={'yes' if entry['markdown'] else 'no'}, "
                          f"json={'yes' if entry['json'] else 'no'}, tables={len(entry['tables'])}")
            elif status == "partial":
                log.warning(f"  \u26a0 {origin.name}: partial success \u2014 {errors}")
            else:
                log.error(f"  \u2717 {origin.name}: failed \u2014 {errors}")

            files_out.append(entry)

        log.info(f"Done in {time.time() - start_time:.1f}s")

    # Clean up temp legacy-xlsx + chunk files. source/ itself is untouched.
    for f in converted_xlsx:
        try:
            f.unlink()
        except Exception as e:
            log.warning(f"  Could not delete temp file {f.name}: {e}")
    shutil.rmtree(work_dir, ignore_errors=True)

    files_out.sort(key=lambda e: e["name"])
    summary = {
        "success": sum(1 for e in files_out if e["status"] == "success"),
        "partial": sum(1 for e in files_out if e["status"] == "partial"),
        "failed": sum(1 for e in files_out if e["status"] in ("failed", "exception")),
    }

    metadata = {
        "project": folder_path.name,
        "source_folder": str(source_dir),
        "files_processed": len(files_out),
        "files": files_out,
        "skipped_unsupported": skipped_unsupported,
        "summary": summary,
        "chunk_threshold": chunk_threshold,
        "last_run": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    }

    with (folder_path / "metadata.json").open("w", encoding="utf-8") as fp:
        json.dump(metadata, fp, ensure_ascii=False, indent=2)

    return metadata
